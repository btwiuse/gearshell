#!/usr/bin/env python3
"""
Build the GearShell · HKUST Round 1 Pitch Deck (PPT).

Output: /Users/gear/Documents/GitHub/gearshell/神奇海螺队 - Round 1 Pitch Deck.pptx
"""

from pathlib import Path
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- design system ----------
BG_DEEP   = RGBColor(0x0B, 0x12, 0x20)   # primary background
BG_PANEL  = RGBColor(0x14, 0x21, 0x3D)   # gradient / panel
BG_CARD   = RGBColor(0x1C, 0x2B, 0x54)   # raised card
LINE      = RGBColor(0x2B, 0x3A, 0x67)   # borders
ACCENT    = RGBColor(0x8A, 0xB4, 0xFF)   # landing-page accent
ACCENT_2  = RGBColor(0xC9, 0xD6, 0xFF)   # soft accent text
MUTED     = RGBColor(0x6B, 0x7A, 0x99)   # captions, footer
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GOOD      = RGBColor(0x6E, 0xE7, 0xB7)   # numbers / good signals
WARN      = RGBColor(0xFF, 0xB4, 0x6E)   # risk highlight
NEG       = RGBColor(0xFF, 0x8A, 0x8A)   # negative trend

FONT_CN = "Microsoft YaHei"
FONT_EN = "Helvetica Neue"

WORKSPACE = Path("/Users/gear/Documents/GitHub/gearshell")
ASSET_SOCIAL   = WORKSPACE / "social.png"
ASSET_ALPHA    = WORKSPACE / "alpha.png"
ASSET_LOGO     = WORKSPACE / "logo-banner-logo.png"
ASSET_GROUP    = WORKSPACE / "group.png"
DEMO_VIDEO     = WORKSPACE / "scripts" / "demo.mp4"

# ---------- helpers ----------

def set_slide_bg(slide, color=BG_DEEP, gradient_to=None):
    bg = slide.background
    fill = bg.fill
    if gradient_to is None:
        fill.solid()
        fill.fore_color.rgb = color
        return
    # gradient background — build the XML directly on cSld/bg
    cSld = slide._element.find(qn("p:cSld"))
    bg_el = cSld.find(qn("p:bg"))
    if bg_el is None:
        bg_el = etree.SubElement(cSld, qn("p:bg"))
    # purge existing children
    for child in list(bg_el):
        bg_el.remove(child)
    bgPr = etree.SubElement(bg_el, qn("p:bgPr"))
    bgPr.set("flip", "none")
    bgPr.set("rotWithShape", "1")
    # gradFill
    gradFill = etree.SubElement(bgPr, qn("a:gradFill"))
    gradFill.set("flip", "none")
    gradFill.set("rotWithShape", "1")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    def add_stop(pos, rgb):
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(int(pos * 1000)))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", rgb)
    add_stop(0.0, "%02X%02X%02X" % (color[0], color[1], color[2]))
    add_stop(1.0, "%02X%02X%02X" % (gradient_to[0], gradient_to[1], gradient_to[2]))
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("angle", "1350000")
    etree.SubElement(gradFill, qn("a:tileRect"))
    # also include bgPr attributes
    etree.SubElement(bgPr, qn("a:effectLst"))


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if False else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    return shp


def add_card(slide, x, y, w, h, fill=BG_CARD, line_color=None, accent_left=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.color.rgb = LINE
    else:
        shp.line.color.rgb = line_color
    shp.line.width = Pt(0.75)
    if accent_left:
        # thin accent stripe down the left side
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(x), Inches(y + 0.05),
                                        Inches(0.06), Inches(h - 0.1))
        stripe.shadow.inherit = False
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = ACCENT
        stripe.line.fill.background()
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_CN,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    if isinstance(text, str):
        runs = [text]
    else:
        runs = text
    first = True
    for chunk in runs:
        if first:
            r = p.add_run() if p.runs else p.add_run()
            first = False
        else:
            r = p.add_run()
        if isinstance(chunk, dict):
            r.text = chunk.get("text", "")
            r.font.size = Pt(chunk.get("size", size))
            r.font.bold = chunk.get("bold", bold)
            r.font.color.rgb = chunk.get("color", color)
            r.font.name = chunk.get("font", font)
        else:
            r.text = chunk
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def add_paragraphs(slide, x, y, w, h, lines, *, size=14, color=WHITE,
                   align=PP_ALIGN.LEFT, font=FONT_CN, line_spacing=1.2,
                   bullet=False):
    """lines: list of plain strings, (text, opts_dict), or (label, body_str)
    for an inline label + body run on the same bullet line."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        bullet_prefix = "• " if bullet else ""
        if isinstance(line, tuple) and len(line) == 2:
            first, second = line
            if isinstance(second, dict):
                text, opts = line
                run_text = bullet_prefix + text
                r = p.add_run()
                r.text = run_text
                r.font.size = Pt(opts.get("size", size))
                r.font.bold = opts.get("bold", False)
                r.font.color.rgb = opts.get("color", color)
                r.font.name = opts.get("font", font)
                continue
            else:
                # (label, body_str) -> label bold, body normal
                r1 = p.add_run()
                r1.text = bullet_prefix + first + "："
                r1.font.size = Pt(size)
                r1.font.bold = True
                r1.font.color.rgb = ACCENT
                r1.font.name = font
                r2 = p.add_run()
                r2.text = second
                r2.font.size = Pt(size)
                r2.font.color.rgb = color
                r2.font.name = font
                continue
        # plain string
        text = line
        if bullet:
            text = "• " + text
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_pill(slide, x, y, text, *, w=None, h=0.32, fill=None, color=ACCENT, size=11):
    if fill is None:
        # translucent feel via solid muted blue
        fill = RGBColor(0x1F, 0x2A, 0x4F)
    if w is None:
        w = max(0.8, 0.18 * len(text) + 0.4)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = FONT_EN
    return shp


def add_section_header(slide, section_num, title, *, kicker=None):
    """Standard section header used on most content slides."""
    if kicker:
        add_text(slide, 0.6, 0.45, 12, 0.4,
                 kicker.upper(), size=11, bold=True, color=ACCENT, font=FONT_EN)
    # section number badge
    add_text(slide, 0.6, 0.75, 0.8, 0.45,
             section_num, size=14, bold=True, color=MUTED, font=FONT_EN)
    add_text(slide, 1.25, 0.7, 11, 0.7,
             title, size=30, bold=True, color=WHITE)
    # accent underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.6), Inches(1.35),
                                       Inches(0.5), Inches(0.06))
    underline.shadow.inherit = False
    underline.fill.solid()
    underline.fill.fore_color.rgb = ACCENT
    underline.line.fill.background()


def add_footer(slide, page_no, total=20):
    add_text(slide, 0.6, 7.1, 6, 0.3,
             "GearShell · 神奇海螺队 · HKUST Million $ · Shanghai · 2026",
             size=9, color=MUTED, font=FONT_EN)
    add_text(slide, 11.5, 7.1, 1.5, 0.3,
             f"{page_no:02d} / {total:02d}", size=9, color=MUTED,
             font=FONT_EN, align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


# ---------- slide builders ----------

def slide_cover(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s, gradient_to=BG_PANEL)

    # accent bar top-left
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0), Inches(0), Inches(0.25), Inches(7.5))
    bar.shadow.inherit = False
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # eyebrow
    add_text(s, 0.9, 0.9, 12, 0.4,
             "HKUST MILLION-DOLLAR ENTREPRENEURSHIP  ·  SHANGHAI  ·  2026",
             size=11, bold=True, color=ACCENT, font=FONT_EN)

    # title block
    add_text(s, 0.9, 1.5, 11.5, 1.6,
             "GearShell", size=72, bold=True, color=WHITE, font=FONT_EN)
    add_text(s, 0.9, 2.7, 11.5, 0.7,
             "浏览器原生 Shell · Web Native Agent Sandbox",
             size=26, bold=True, color=ACCENT, font=FONT_CN)
    add_text(s, 0.9, 3.45, 11.5, 0.6,
             "让 AI Agent 像网页一样开箱即用",
             size=18, color=ACCENT_2, font=FONT_CN)

    # one-liner
    add_text(s, 0.9, 4.4, 11.5, 0.5,
             "A kernel. A shell. A terminal. A browser. An AI assistant.  ·  All in one tab.",
             size=14, color=WHITE, font=FONT_EN)

    # bottom info block
    add_rect(s, 0.9, 5.5, 11.5, 1.2, fill=BG_CARD, line_color=LINE)
    add_text(s, 1.1, 5.6, 4.5, 0.4, "团队 / Team", size=10, bold=True,
             color=ACCENT, font=FONT_EN)
    add_text(s, 1.1, 5.9, 5, 0.4, "神奇海螺队", size=18, bold=True, color=WHITE)
    add_text(s, 1.1, 6.3, 6, 0.3, "邹航标 (Hangbiao Zou) · Founder & 核心开发者",
             size=11, color=ACCENT_2)

    add_text(s, 6.5, 5.6, 4.5, 0.4, "赛道 / Track", size=10, bold=True,
             color=ACCENT, font=FONT_EN)
    add_text(s, 6.5, 5.9, 5, 0.4, "人工智能 / 开发者基础设施", size=18, bold=True,
             color=WHITE)
    add_text(s, 6.5, 6.3, 6, 0.3, "初创组 · 估值 < 2 亿 RMB", size=11, color=ACCENT_2)

    # corner mark
    add_text(s, 11.0, 7.1, 2, 0.3, "Round 1 Pitch Deck · 2026",
             size=9, color=MUTED, font=FONT_EN, align=PP_ALIGN.RIGHT)
    add_notes(s, "Cover · 自我介绍 · 15s 之内讲清楚三件事：是什么、为谁、为什么是现在。")


def slide_agenda(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_text(s, 0.6, 0.7, 12, 0.5,
             "AGENDA", size=12, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 0.6, 1.05, 12, 1,
             "我们要在 7 分钟内讲清楚什么", size=34, bold=True, color=WHITE)

    items = [
        ("01", "The Gap", "AI Agent 分发的「最后一公里」"),
        ("02", "Solution", "浏览器原生 Shell · Zero Install · Local-first"),
        ("03", "Architecture", "WASM 微内核 / VFS / 终端 / Userland 四层"),
        ("04", "Why Us", "自举验证 + ChromiumOS 一线工程经验"),
        ("05", "Market & Customers", "TAM / SAM / SOM · ToB / ToD / ToC"),
        ("06", "Business & GTM", "四条收入线 · 三步入场"),
        ("07", "Roadmap & Ask", "里程碑 · 18 个月 300 万融资计划"),
    ]
    col_w = 5.7
    for i, (num, head, sub) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.6 + col * (col_w + 0.3)
        y = 2.6 + row * 1.1
        add_card(s, x, y, col_w, 0.95)
        add_text(s, x + 0.3, y + 0.12, 0.8, 0.4,
                 num, size=22, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, x + 1.2, y + 0.1, 4.5, 0.4,
                 head, size=18, bold=True, color=WHITE, font=FONT_EN)
        add_text(s, x + 1.2, y + 0.5, 4.4, 0.4,
                 sub, size=12, color=ACCENT_2)
    add_footer(s, 2)
    add_notes(s, "Agenda · 路线图，2 分钟定位（前 4 项），3 分钟商业（后 3 项），2 分钟 Q&A 留白。")


def slide_tldr(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s, gradient_to=BG_PANEL)
    add_text(s, 0.6, 0.7, 12, 0.5, "TL;DR", size=12, bold=True,
             color=ACCENT, font=FONT_EN)
    add_text(s, 0.6, 1.05, 12, 0.8, "一页执行摘要", size=34, bold=True, color=WHITE)

    # big quote
    add_card(s, 0.6, 2.1, 12.1, 1.4, accent_left=True)
    add_text(s, 0.95, 2.3, 11.7, 0.4,
             "我们相信：浏览器就是下一个操作系统，URL 就是下一个二进制。",
             size=20, bold=True, color=WHITE, font=FONT_CN)
    add_text(s, 0.95, 2.85, 11.7, 0.5,
             "AI Agent 不应该需要「安装」——它应该像网页一样被分发，在用户的浏览器里直接运行。",
             size=14, color=ACCENT_2, font=FONT_CN)

    # three pillars
    pillars = [
        ("WHAT",
         "浏览器原生 Shell",
         "把 Kernel · Shell · Terminal · 浏览器 · AI 助手装进一个 Tab"),
        ("FOR WHOM",
         "Agent 厂商 + 企业 IT + 开发者",
         "解决「不会装、不敢用、生态割裂」三大 Agent 分发痛点"),
        ("WHY NOW",
         "WebAssembly + 浏览器本地 AI",
         "WASM 开放标准成熟 + WebGPU 让本地 27B 模型可在浏览器内推理"),
    ]
    for i, (tag, head, sub) in enumerate(pillars):
        x = 0.6 + i * 4.07
        add_card(s, x, 3.85, 3.87, 1.6)
        add_text(s, x + 0.3, 3.95, 3.5, 0.3,
                 tag, size=10, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, x + 0.3, 4.25, 3.5, 0.4,
                 head, size=18, bold=True, color=WHITE, font=FONT_CN)
        add_text(s, x + 0.3, 4.75, 3.5, 1.0,
                 sub, size=12, color=ACCENT_2, font=FONT_CN, line_spacing=1.3)

    # bottom ask strip
    ask = add_card(s, 0.6, 5.7, 12.1, 1.0, fill=BG_CARD, line_color=ACCENT)
    add_text(s, 0.95, 5.85, 4, 0.3,
             "本轮 ASKING", size=10, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 0.95, 6.15, 5.5, 0.4,
             "300 万 RMB · 18 个月", size=22, bold=True, color=WHITE, font=FONT_EN)
    add_text(s, 7.0, 5.85, 5.5, 0.3,
             "USE OF FUNDS", size=10, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 7.0, 6.15, 5.5, 0.4,
             "60% 研发 · 20% 市场 · 15% 团队 · 5% 基础设施",
             size=14, color=WHITE, font=FONT_CN)

    add_footer(s, 3)
    add_notes(s, "TL;DR · 评委看完这一页就能复述出 GearShell 是什么、本轮要什么。3 大投资亮点先抛。")


def slide_problem(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "01", "AI Agent 分发的「最后一公里」",
                       kicker="THE GAP")

    # three pain cards
    cards = [
        ("不会装", "Install Friction",
         "技术小白因安装复杂、系统不兼容而跳出率高；Agent 厂商获客漏斗在「安装」环节流失 > 60%。",
         "跳出率", "60%+"),
        ("不敢用", "Trust & Compliance",
         "企业 IT 担忧 Agent 权限滥用、数据泄露；现有本地工具缺乏默认隔离与可观测性。",
         "企业内测", "≤ 3 个月"),
        ("生态割裂", "Ecosystem Lock-in",
         "主流沙盒仅支持单一语言（Node.js / Python）；iPad 等平台覆盖差，跨设备体验断裂。",
         "语言支持", "1 种"),
    ]
    for i, (head, en, body, kpi_l, kpi_v) in enumerate(cards):
        x = 0.6 + i * 4.07
        add_card(s, x, 1.7, 3.87, 4.0, accent_left=True)
        add_text(s, x + 0.3, 1.9, 3.5, 0.5,
                 head, size=24, bold=True, color=WHITE, font=FONT_CN)
        add_text(s, x + 0.3, 2.45, 3.5, 0.3,
                 en, size=10, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, x + 0.3, 2.8, 3.5, 1.4,
                 body, size=12, color=ACCENT_2, font=FONT_CN, line_spacing=1.4)
        add_rect(s, x + 0.3, 4.7, 3.4, 0.85, fill=BG_PANEL, line_color=LINE)
        add_text(s, x + 0.45, 4.78, 1.5, 0.3,
                 kpi_l.upper(), size=9, bold=True, color=MUTED, font=FONT_EN)
        add_text(s, x + 0.45, 5.05, 3.2, 0.45,
                 kpi_v, size=22, bold=True, color=ACCENT, font=FONT_EN)

    # bottom callout
    add_text(s, 0.6, 6.0, 12.1, 0.6,
             "「又一个 Agent 装完就卸，xxx 文件夹还得手动清」——X 用户 @Zachary_haha",
             size=14, color=WARN, font=FONT_CN)
    add_text(s, 0.6, 6.45, 12.1, 0.4,
             "Zachary_haha · 2026  ·  Source: x.com/Zachary_haha/status/2084644286042198287",
             size=9, color=MUTED, font=FONT_EN)
    add_footer(s, 4)
    add_notes(s, "Problem · 三个真实痛点 + 量化影响 + 一句来自真实用户的吐槽佐证。")


def slide_voice_from_field(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "01'", "市场信号：从业者已经在选「浏览器」",
                       kicker="VOICE FROM THE FIELD")

    # left: david cramer poll card
    add_card(s, 0.6, 1.7, 6.0, 5.0, accent_left=True)
    add_text(s, 0.85, 1.85, 5.5, 0.35,
             "DAVID CRAMER · @zeeg · 2026-06-27", size=10, bold=True,
             color=ACCENT, font=FONT_EN)
    add_text(s, 0.85, 2.15, 5.5, 0.7,
             "「If you're building agents... where are you hosting them?」",
             size=14, bold=True, color=WHITE, font=FONT_EN)

    # poll bars
    poll = [
        ("Cloudflare", 43.6),
        ("AWS / GCP / Azure", 24.7),
        ("Vercel", 11.2),
        ("Other", 20.5),
    ]
    bar_x = 0.85
    bar_y = 3.1
    bar_max = 4.7
    for label, pct in poll:
        add_text(s, bar_x, bar_y, 5.5, 0.3,
                 label, size=11, color=WHITE)
        add_rect(s, bar_x, bar_y + 0.3, bar_max, 0.22, fill=BG_PANEL, line_color=LINE)
        add_rect(s, bar_x, bar_y + 0.3, bar_max * (pct / 50.0), 0.22,
                 fill=ACCENT, line_color=None)
        add_text(s, bar_x + bar_max + 0.05, bar_y + 0.27, 0.7, 0.3,
                 f"{pct:.1f}%", size=11, bold=True, color=ACCENT, font=FONT_EN)
        bar_y += 0.55
    add_text(s, 0.85, 5.7, 5.5, 0.3,
             "2,101 votes · Final results",
             size=10, color=MUTED, font=FONT_EN)
    add_text(s, 0.85, 6.05, 5.5, 0.4,
             "过半数 Agent 团队在云上跑；浏览器是 Next。",
             size=12, color=ACCENT_2, font=FONT_CN)

    # right: jeff lindsay reply
    add_card(s, 6.85, 1.7, 6.0, 5.0, accent_left=True, line_color=ACCENT)
    add_text(s, 7.1, 1.85, 5.5, 0.35,
             "JEFF LINDSAY · @progrium · 6:43 AM · 2026-06-27",
             size=10, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 7.1, 2.15, 5.5, 0.3,
             "Reply to the thread:", size=10, color=MUTED, font=FONT_EN)
    add_text(s, 7.1, 2.55, 5.5, 1.6,
             "\u201Ci dont know what this means but im pretty sure my answer is in the browser\u201D",
             size=15, bold=True, color=WHITE, font=FONT_EN, line_spacing=1.3)

    add_text(s, 7.1, 4.4, 5.5, 0.3,
             "WHY THIS MATTERS", size=10, bold=True, color=ACCENT, font=FONT_EN)
    add_paragraphs(s, 7.1, 4.7, 5.5, 2.0, [
        "Jeff Lindsay：Docker / Webhooks 推动者，国际基础设施社区 KOL",
        "在 2,100 名 Agent 工程师的投票底下说「我的答案是浏览器」",
        "对应 GearShell：浏览器原生 Shell = Agent 的默认 runtime",
    ], size=12, color=ACCENT_2, line_spacing=1.4, bullet=True)

    add_footer(s, 5)
    add_notes(s, "Voice · 用真实社区投票和 Jeff Lindsay 的回复作为外部背书，避免「我们觉得」。" \
               "源：x.com/zeeg/status/2070591092471558567 与 x.com/progrium/status/2070639004761145654。")


def slide_solution(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "02", "Solution：浏览器原生 Shell",
                       kicker="SOLUTION")

    # hero quote card
    add_card(s, 0.6, 1.7, 12.1, 1.3, accent_left=True)
    add_text(s, 0.95, 1.85, 11.7, 0.4,
             "ONE-LINER", size=10, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 0.95, 2.2, 11.7, 0.7,
             "一层运行在主流浏览器之上的轻量虚拟操作环境——",
             size=18, color=WHITE, font=FONT_CN)
    add_text(s, 0.95, 2.6, 11.7, 0.4,
             "为企业、开发者和终端用户提供受控的进程、文件系统、网络与终端能力。",
             size=14, color=ACCENT_2, font=FONT_CN)

    # three value pillars
    pillars = [
        ("极致轻量 · 免安装",
         "Zero Install",
         "核心组件仅依赖 CDN 静态资源，可嵌入任意 Web 应用，零服务器资源占用。"),
        ("Local-first · 默认安全",
         "Browser = Root of Trust",
         "默认 VFS 隔离、进程受控、可观测；数据不出浏览器，Agent 像网页一样被分发。"),
        ("多语言生态兼容",
         "WASI · Multi-language",
         "基于 WebAssembly 标准，C++ / Rust / Go / TypeScript 同等运行，远超单一生态。"),
    ]
    for i, (head, en, body) in enumerate(pillars):
        x = 0.6 + i * 4.07
        add_card(s, x, 3.25, 3.87, 2.6, accent_left=True)
        add_text(s, x + 0.3, 3.4, 3.5, 0.5,
                 head, size=16, bold=True, color=WHITE, font=FONT_CN)
        add_text(s, x + 0.3, 3.85, 3.5, 0.3,
                 en, size=10, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, x + 0.3, 4.2, 3.5, 1.6,
                 body, size=12, color=ACCENT_2, font=FONT_CN, line_spacing=1.4)

    # differentiators strip
    add_text(s, 0.6, 6.1, 12.1, 0.4,
             "对比维度差异化（详见后页 Competitive 表格）：",
             size=11, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 0.6, 6.45, 12.1, 0.5,
             "浏览器即开即用 · 完整 PTY 终端 · 自研微内核 + VFS · 规划中 Agent 应用商店",
             size=13, color=WHITE, font=FONT_CN)
    add_footer(s, 6)
    add_notes(s, "Solution · 用一句话定义产品；用三个支柱展示差异化。避免堆 feature list。")


def slide_demo(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "02'", "Live Demo · 打开一个 Tab，拿到一个 OS",
                       kicker="PRODUCT DEMO")

    # left: terminal card (mimic landing demo card)
    add_card(s, 0.6, 1.7, 7.0, 4.9, fill=BG_PANEL)
    # title bar
    add_rect(s, 0.6, 1.7, 7.0, 0.5, fill=BG_CARD, line_color=None)
    for i, c in enumerate([(0xFF, 0x6B, 0x6B), (0xFF, 0xC1, 0x6E), (0x6E, 0xE7, 0xB7)]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(0.78 + i * 0.28), Inches(1.85),
                               Inches(0.18), Inches(0.18))
        d.shadow.inherit = False
        d.fill.solid()
        d.fill.fore_color.rgb = RGBColor(*c)
        d.line.fill.background()
    add_text(s, 2.0, 1.78, 5.5, 0.4,
             "gear@gear: ~", size=11, bold=True,
             color=ACCENT_2, font=FONT_EN)

    # terminal body
    body = [
        ("$ gear init my-app", GOOD),
        ("  ✓ pulling linux userspace (47 MB)", ACCENT_2),
        ("  ✓ mounting vfs at /home/gear", ACCENT_2),
        ("  ✓ ready in 312 ms", GOOD),
        ("", WHITE),
        ("$ cd my-app && git clone github.com/me/repo", GOOD),
        ("Cloning into 'repo'...", ACCENT_2),
        ("✓ done. 142 files, 3 branches.", GOOD),
        ("", WHITE),
        ("$ claude", GOOD),
        ("> Hi, I am Claude Code. Working in /home/gear/my-app/repo.", ACCENT),
        ("> What would you like to build?", ACCENT),
        ("$ ▌", WHITE),
    ]
    tb = s.shapes.add_textbox(Inches(0.85), Inches(2.35),
                              Inches(6.5), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05)
    for i, (line, c) in enumerate(body):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.name = "Menlo"
        r.font.color.rgb = c

    add_text(s, 0.6, 6.65, 7.0, 0.3,
             "Real terminal, real Claude Code, zero installs.",
             size=10, bold=True, color=MUTED, font=FONT_EN)

    # right: video placeholder + checklist
    add_card(s, 7.85, 1.7, 4.8, 4.9, accent_left=True)
    add_text(s, 8.1, 1.85, 4.4, 0.35,
             "PRODUCT VIDEO · 嵌入位", size=10, bold=True,
             color=ACCENT, font=FONT_EN)
    add_text(s, 8.1, 2.2, 4.4, 0.5,
             "▶  产品演示短视频位",
             size=16, bold=True, color=WHITE, font=FONT_CN)

    # placeholder area
    ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(8.1), Inches(2.75),
                            Inches(4.3), Inches(2.3))
    ph.shadow.inherit = False
    ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0x10, 0x1A, 0x33)
    ph.line.color.rgb = LINE
    ph.line.width = Pt(1.0)
    add_text(s, 8.1, 3.2, 4.3, 0.4,
             "DEMO_VIDEO_PLACEHOLDER",
             size=11, bold=True, color=ACCENT, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_text(s, 8.1, 3.55, 4.3, 0.4,
             "scripts/demo.mp4", size=10, color=MUTED, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_text(s, 8.1, 3.85, 4.3, 0.4,
             "正式路演前替换为 60–90s 产品录屏",
             size=11, color=ACCENT_2, font=FONT_CN, align=PP_ALIGN.CENTER)

    # features strip inside card
    add_text(s, 8.1, 5.2, 4.4, 0.3,
             "WHAT YOU CAN DO IN 60 SECONDS", size=10, bold=True,
             color=ACCENT, font=FONT_EN)
    add_paragraphs(s, 8.1, 5.5, 4.4, 1.1, [
        "打开 gear.sh → 输入 `gear init` → 拿到真实 Linux 用户态",
        "克隆 repo → 运行 Claude Code → 状态跨设备持久化",
        "VFS 隔离 → 关闭再开 Tab → 工作目录还在",
    ], size=11, color=WHITE, line_spacing=1.4, bullet=True)

    add_footer(s, 7)
    add_notes(s,
        "Demo · 现场最好跑一次 gear.sh。PPT 中右侧为视频占位，"
        "替换路径：scripts/demo.mp4（正式路演前用 QuickTime 录 60-90s 录屏，"
        "保存到 scripts/demo.mp4，重新跑 python3 scripts/build_pitch_deck.py 即嵌入）。"
    )


def slide_architecture(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "03", "技术架构：四层自研体系",
                       kicker="ARCHITECTURE")

    # layer stack
    layers = [
        ("Userland",
         "Shell · Agent · 工具集 · 包管理器 · Registry",
         "已可运行：git · docker · kubectl · esbuild · tsc · go · claude code · crush"),
        ("终端与网络",
         "完整 PTY 终端模拟 · 受控可观测网络能力",
         "xterm.js 驱动，CLI Agent 无修改运行"),
        ("虚拟文件系统 (VFS)",
         "独立受控文件工作区 · 精细权限 · 本地/远程挂载",
         "支持 OPFS + 远程仓库导入，状态可持久到 URL"),
        ("WASM 微内核",
         "WebAssembly 多进程管理 + 资源隔离",
         "Linux syscall · 真实进程 · 真实文件系统 · 真实网络"),
    ]
    base_y = 1.7
    layer_h = 1.1
    for i, (head, mid, sub) in enumerate(layers):
        y = base_y + i * (layer_h + 0.18)
        add_card(s, 0.6, y, 12.1, layer_h, accent_left=(i == 3))
        # tag
        add_text(s, 0.95, y + 0.15, 1.3, 0.4,
                 f"L{4-i}", size=11, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, 2.4, y + 0.1, 5, 0.4,
                 head, size=18, bold=True, color=WHITE, font=FONT_CN)
        add_text(s, 2.4, y + 0.55, 6.5, 0.5,
                 mid, size=12, color=ACCENT_2, font=FONT_CN)
        add_text(s, 8.7, y + 0.3, 4, 0.6,
                 sub, size=11, color=WHITE, font=FONT_CN, line_spacing=1.3)

    # self-host callout
    add_text(s, 0.6, 6.55, 12.1, 0.4,
             "自举验证 (Self-host)：GearShell 已经在 GearShell 里构建 GearShell。",
             size=13, bold=True, color=GOOD, font=FONT_CN)
    add_text(s, 0.6, 6.9, 12.1, 0.3,
             "意味着生态可以从一个 Tab 内自我繁殖——这是后续 Agent 应用商店的工程前提。",
             size=11, color=ACCENT_2, font=FONT_CN)
    add_footer(s, 8)
    add_notes(s, "Architecture · 四层从下到上是依赖关系。最下面 WASM 微内核是最深的护城河。" \
               "「自举验证」是要点：意味着这个平台可以跑自己的开发流程。")


def slide_why_us(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "04", "Why Us：这条赛道为什么是我们",
                       kicker="WHY US")

    # left: founder track record
    add_card(s, 0.6, 1.7, 6.0, 5.0, accent_left=True)
    add_text(s, 0.85, 1.85, 5.5, 0.35,
             "FOUNDER · 邹航标 (Hangbiao Zou)",
             size=11, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 0.85, 2.2, 5.5, 0.4,
             "WebAssembly / 浏览器架构深度实践者",
             size=15, bold=True, color=WHITE, font=FONT_CN)
    add_paragraphs(s, 0.85, 2.7, 5.5, 3.7, [
        ("FydeOS · 国内首家 ChromiumOS 团队", {"bold": True}),
        "从底层运行时到应用层全栈工程经验（系统 / 安全 / 前端）",
        ("WebAssembly 生态长期贡献者", {"bold": True}),
        "WASI 标准持续跟进 · 多语言 sandbox 实践经验",
        ("国际开发者社区影响", {"bold": True}),
        "RustConf China / EthShanghai / EthShenzhen 演讲嘉宾",
        ("AI Agent 运行时与开源社区经验", {"bold": True}),
        "Hacker News · 黑客松 · 与 Jeff Lindsay 等常态交流",
    ], size=11, color=ACCENT_2, line_spacing=1.4, bullet=True)

    # right: 4 reasons
    add_text(s, 6.85, 1.7, 6.0, 0.4,
             "为什么是这个团队做这件事",
             size=16, bold=True, color=WHITE, font=FONT_CN)
    reasons = [
        ("1. 唯一全栈经验",
         "从 ChromiumOS 内核到 Web 前端，没有任何团队比我们更懂「浏览器当 OS」这件事。"),
        ("2. 自举已经发生",
         "GearShell 已经在 GearShell 里构建自己——这是别人用嘴说、我们用工程证明的护城河。"),
        ("3. 社区到位",
         "Bonsai 27B WebGPU · Wagi Dog 桌面宠物 · FydeOS 安全团队 — 子模块已经构成产品矩阵。"),
        ("4. 已验证场景",
         "Infra 团队已把 GearShell 嵌进文档（参考 @Devin K 案例），无需 Docker / npm / Slack 求助。"),
    ]
    for i, (h, b) in enumerate(reasons):
        y = 2.2 + i * 1.13
        add_card(s, 6.85, y, 6.0, 1.0, fill=BG_PANEL)
        add_text(s, 7.05, y + 0.1, 5.8, 0.35,
                 h, size=13, bold=True, color=ACCENT, font=FONT_CN)
        add_text(s, 7.05, y + 0.45, 5.8, 0.55,
                 b, size=11, color=WHITE, font=FONT_CN, line_spacing=1.35)

    add_footer(s, 9)
    add_notes(s, "Why Us · 创始人在 ChromiumOS / FydeOS 干过，浏览器架构是底子。" \
               "「自举」是工程能力证明，不是 PPT 写法。")


def slide_traction(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "04'", "Traction & Validation",
                       kicker="TRACTION")

    # KPI cards
    kpis = [
        ("100% Buildless",
         "0 npm install · 0 webpack · 0 vite",
         "importmap + CDN 即可分发整套 Shell"),
        ("4 子模块",
         "Bonsai · Browser.js · Isolation · Workbench",
         "已在 gear.sh 同时跑通，全部 git submodule"),
        ("自举验证",
         "GearShell ⊂ GearShell",
         "在浏览器内构建、编译、运行自身"),
        ("27B WebGPU",
         "本地 LLM 推理",
         "Bonsai 27B · 3.8 GB · 浏览器内推理 · 不出本机"),
    ]
    for i, (head, mid, sub) in enumerate(kpis):
        x = 0.6 + i * 3.1
        add_card(s, x, 1.7, 2.9, 1.9, accent_left=True)
        add_text(s, x + 0.2, 1.85, 2.6, 0.4,
                 head, size=18, bold=True, color=WHITE, font=FONT_EN)
        add_text(s, x + 0.2, 2.35, 2.6, 0.6,
                 mid, size=11, color=ACCENT, font=FONT_CN)
        add_text(s, x + 0.2, 2.95, 2.6, 0.6,
                 sub, size=10, color=MUTED, font=FONT_CN, line_spacing=1.3)

    # bottom: timeline + early validation
    add_text(s, 0.6, 3.85, 12, 0.4,
             "DEVELOPMENT TIMELINE · 发展大事记",
             size=11, bold=True, color=ACCENT, font=FONT_EN)
    timeline = [
        ("2024 Q4", "WASM 微内核 PoC", "在浏览器内跑通第一个 Linux syscall"),
        ("2025 Q2", "VFS + 多语言运行时", "git / docker / kubectl 全量移植"),
        ("2025 Q4", "自举 + Workbench", "在 GearShell 里编译 GearShell · 内嵌 VS Code-OSS"),
        ("2026 Q1", "Bonsai 27B + WebGPU", "本地 27B 模型浏览器推理，无服务端"),
        ("2026 Q2", "gear.sh 上线", "公开 SaaS 入口 + Landing Page 改版"),
        ("2026 Q3", "HKUST 参赛", "面向上海科创生态提交 Round 1 材料"),
    ]
    base_x = 0.6
    col_w = 12.1 / 6
    for i, (when, head, body) in enumerate(timeline):
        x = base_x + i * col_w
        # dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x + col_w/2 - 0.08), Inches(4.3),
                                 Inches(0.16), Inches(0.16))
        dot.shadow.inherit = False
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(s, x + 0.05, 4.5, col_w - 0.1, 0.3,
                 when, size=10, bold=True, color=ACCENT, font=FONT_EN,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, 4.8, col_w - 0.1, 0.5,
                 head, size=11, bold=True, color=WHITE, font=FONT_CN,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)
        add_text(s, x + 0.05, 5.3, col_w - 0.1, 0.6,
                 body, size=9, color=ACCENT_2, font=FONT_CN,
                 align=PP_ALIGN.CENTER, line_spacing=1.25)

    # bottom: validation quote
    add_card(s, 0.6, 6.1, 12.1, 0.85, fill=BG_PANEL)
    add_text(s, 0.85, 6.2, 11.7, 0.3,
             "EARLY VALIDATION",
             size=10, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 0.85, 6.45, 11.7, 0.4,
             "「We embed a GearShell tab in our docs. Customers run the example in one click — "
             "no Docker, no npm install, no Slack message asking us why their node version is wrong.」",
             size=12, color=WHITE, font=FONT_EN, line_spacing=1.3)
    add_text(s, 0.85, 6.75, 11.7, 0.3,
             "— Devin K. · Infra Lead, Series A · 2026",
             size=10, color=MUTED, font=FONT_EN)

    add_footer(s, 10)
    add_notes(s, "Traction · KPI + 时间线 + 客户证言。指标避免虚高，用「已经发生」的事实陈述。")


def slide_market(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "05", "市场机遇：浏览器 + Agent 的交叉点",
                       kicker="MARKET")

    # TAM/SAM/SOM bars
    sizes = [
        ("TAM", "全球浏览器端 AI Agent 沙盒", "USD 47B", "2030E · CAGR ~45%"),
        ("SAM", "中国开发者 + 企业 AI 工具市场", "RMB 260亿", "2026E · 增速全球领先"),
        ("SOM", "首批 Agent 厂商 + 大客户", "RMB 3-5亿", "12-18 个月内可达 ARR"),
    ]
    add_text(s, 0.6, 1.65, 12, 0.4,
             "MARKET SIZE · 第三方测算",
             size=11, bold=True, color=ACCENT, font=FONT_EN)

    base_y = 2.1
    row_h = 1.3
    for i, (tag, head, size_v, sub) in enumerate(sizes):
        y = base_y + i * row_h
        add_card(s, 0.6, y, 12.1, 1.05)
        add_text(s, 0.85, y + 0.15, 1.0, 0.4,
                 tag, size=18, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, 1.95, y + 0.15, 4.5, 0.4,
                 head, size=15, bold=True, color=WHITE, font=FONT_CN)
        # size bar
        bar_total = 4.5
        ratios = [1.0, 0.55, 0.18]  # TAM full, SAM ~55%, SOM ~18%
        add_rect(s, 6.5, y + 0.25, bar_total, 0.4, fill=BG_PANEL, line_color=LINE)
        add_rect(s, 6.5, y + 0.25, bar_total * ratios[i], 0.4,
                 fill=ACCENT, line_color=None)
        add_text(s, 11.3, y + 0.2, 1.3, 0.4,
                 size_v, size=14, bold=True, color=ACCENT, font=FONT_EN,
                 align=PP_ALIGN.RIGHT)
        add_text(s, 1.95, y + 0.6, 9.5, 0.35,
                 sub, size=10, color=MUTED, font=FONT_EN)

    # drivers
    add_text(s, 0.6, 6.1, 12, 0.4,
             "THREE GROWTH DRIVERS",
             size=11, bold=True, color=ACCENT, font=FONT_EN)
    drivers = [
        ("技术驱动", "WebAssembly 开放标准快速演进，Agent 得以高安全、零配置运行于浏览器"),
        ("政策驱动", "GDPR · 欧盟 AI 法案 · 中国算法备案推动平台强化隔离与数据合规"),
        ("需求驱动", "企业对「免安装、零配置、强隔离」的分发层需求与 AI Agent 普及正相关"),
    ]
    for i, (h, b) in enumerate(drivers):
        x = 0.6 + i * 4.07
        add_card(s, x, 6.5, 3.87, 0.55, fill=BG_PANEL)
        add_text(s, x + 0.15, 6.55, 3.5, 0.3,
                 h, size=11, bold=True, color=ACCENT, font=FONT_CN)
        add_text(s, x + 0.15, 6.78, 3.6, 0.3,
                 b, size=9, color=WHITE, font=FONT_CN)

    add_footer(s, 11)
    add_notes(s, "Market · 用 TAM/SAM/SOM 阶梯让评委快速定位市场大小；drivers 解释为什么是现在。")


def slide_customers(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "05'", "目标客户：三类核心用户群",
                       kicker="CUSTOMERS")

    segments = [
        ("ToB · 企业 IT / 采购",
         "迁移成本高 · 合规压力大 · 权限滥用风险",
         "流程自动化集成 · 数据处理 · 安全评估试用 · 内部流程开发协作",
         ["金融/政企内网", "零信任落地", "Agent 内测沙盒", "SDK 私有化"]),
        ("ToD · 开发者基础设施",
         "工具生态繁杂 · 安全隔离难 · 分发转化率低",
         "插件/模块分发 · API 接入 · 社区共享 · Agent 试用 DEMO",
         ["开源项目作者", "ISV 工具厂商", "黑客松团队", "AI 工程师"]),
        ("ToC · 终端创新者用户",
         "不会安装 · 害怕破坏系统或泄露隐私",
         "在线写作辅助 · 代码生成 · AI 应用试水 · 临时自动化脚本",
         ["学生 / 自学者", "iPad / 平板用户", "跨设备远程办公", "非工程岗尝鲜"]),
    ]
    for i, (head, pain, scenarios, chips) in enumerate(segments):
        x = 0.6 + i * 4.07
        add_card(s, x, 1.7, 3.87, 4.95, accent_left=True)
        add_text(s, x + 0.3, 1.85, 3.5, 0.5,
                 head, size=16, bold=True, color=WHITE, font=FONT_CN)
        add_text(s, x + 0.3, 2.4, 3.5, 0.4,
                 "核心痛点", size=10, bold=True, color=ACCENT, font=FONT_CN)
        add_text(s, x + 0.3, 2.7, 3.5, 0.7,
                 pain, size=11, color=ACCENT_2, font=FONT_CN, line_spacing=1.3)
        add_text(s, x + 0.3, 3.5, 3.5, 0.4,
                 "典型场景", size=10, bold=True, color=ACCENT, font=FONT_CN)
        add_text(s, x + 0.3, 3.8, 3.5, 1.0,
                 scenarios, size=11, color=WHITE, font=FONT_CN, line_spacing=1.3)

        # chips at bottom
        chip_y = 5.4
        for chip in chips:
            chip_w = max(0.9, 0.18 * len(chip) + 0.3)
            add_pill(s, x + 0.3, chip_y, chip, w=chip_w, h=0.32, size=10)
            chip_y += 0.45

    # bottom note
    add_text(s, 0.6, 6.85, 12.1, 0.3,
             "三类用户群共享同一个产品形态——只是 license / SLA / 集成方式不同，研发成本边际递减。",
             size=11, color=GOOD, font=FONT_CN)
    add_footer(s, 12)
    add_notes(s, "Customers · 三个 segment 的痛点都不一样，但都能被同一个产品解决。")


def slide_business(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "06", "商业模式：四条收入线",
                       kicker="BUSINESS MODEL")

    lines = [
        ("SaaS 订阅 / 席位制",
         "600–1000 RMB / 席位 / 年",
         "按月/年/席位授权 · 中小团队与企业轻量办公",
         "Recurring"),
        ("插件与 Agent 分成",
         "15–30% 分成",
         "开发者经 Registry 销售插件或 Agent，平台抽佣",
         "Marketplace"),
        ("本地部署 / 定制服务",
         "5–20 万 RMB / 单",
         "大客户私有化部署与深度场景定制 · 高毛利",
         "Enterprise"),
        ("行业增值服务",
         "按效果计费 (GaaS)",
         "集成运维 · 合规顾问 · 场景定制咨询",
         "Outcome-based"),
    ]
    for i, (head, price, body, tag) in enumerate(lines):
        x = 0.6 + i * 3.07
        add_card(s, x, 1.7, 2.87, 3.6, accent_left=True)
        add_pill(s, x + 2.05, 1.85, tag, w=1.05, h=0.3, size=9)
        add_text(s, x + 0.25, 1.85, 2.4, 0.4,
                 head, size=15, bold=True, color=WHITE, font=FONT_CN)
        add_text(s, x + 0.25, 2.4, 2.6, 0.6,
                 price, size=18, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, x + 0.25, 3.1, 2.5, 1.6,
                 body, size=11, color=ACCENT_2, font=FONT_CN, line_spacing=1.4)

    # mix shift strip
    add_text(s, 0.6, 5.55, 12, 0.4,
             "REVENUE MIX · 12-18 个月演进",
             size=11, bold=True, color=ACCENT, font=FONT_EN)
    add_card(s, 0.6, 5.95, 12.1, 1.1, fill=BG_PANEL)
    stages = [
        ("2026 H2", "70% 定制 · 25% 订阅 · 5% 分成"),
        ("2027",     "45% 订阅 · 35% 定制 · 15% 分成 · 5% 增值"),
        ("2028",     "55% 订阅 · 20% 分成 · 15% 定制 · 10% 增值"),
        ("2029+",    "45% 订阅 · 35% 分成 · 10% 定制 · 10% 增值"),
    ]
    col_w = 12.1 / 4
    for i, (when, mix) in enumerate(stages):
        x = 0.6 + i * col_w
        add_text(s, x + 0.1, 6.05, col_w - 0.2, 0.3,
                 when, size=11, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, x + 0.1, 6.4, col_w - 0.2, 0.5,
                 mix, size=12, color=WHITE, font=FONT_CN, line_spacing=1.3)

    add_footer(s, 13)
    add_notes(s, "Business · 四条线抵抗单一收入风险；12-18 个月内订阅 + 分成占比提升。")


def slide_gtm(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "06'", "Go-to-Market：三步入场",
                       kicker="GO-TO-MARKET")

    steps = [
        ("01 · 开发者先行",
         "0–6 个月",
         "公共 SDK + 文档开源 + importmap 一行嵌入",
         ("开发者社区", "RustConf / HN / 中文开发者大会"),
         ("关键动作", "黑客松 · 种子计划 · KOL 邀请"),
         ("成功指标", "SDK 周下载 · 嵌入示例数量")),
        ("02 · 行业深耕",
         "6–12 个月",
         "针对细分场景（金融 / 教育 / 政企）打造垂直 DEMO",
         ("渠道", "ISV 联盟 · 云市场集成"),
         ("关键动作", "联合解决方案 · 行业认证材料"),
         ("成功指标", "首批付费客户 · 行业案例")),
        ("03 · 企业采购",
         "12–18 个月",
         "补齐合规与安全认证 · 推行标准试用流程",
         ("销售", "直销顾问 + 客户成功经理"),
         ("关键动作", "私有化部署 · SLA / 等保 2.0"),
         ("成功指标", "ARR · 大客户续约率")),
    ]
    for i, (head, when, body, *rows) in enumerate(steps):
        x = 0.6 + i * 4.07
        add_card(s, x, 1.7, 3.87, 5.0, accent_left=True)
        add_text(s, x + 0.3, 1.85, 3.5, 0.4,
                 when, size=10, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, x + 0.3, 2.2, 3.5, 0.5,
                 head, size=16, bold=True, color=WHITE, font=FONT_CN)
        add_text(s, x + 0.3, 2.7, 3.5, 0.7,
                 body, size=11, color=ACCENT_2, font=FONT_CN, line_spacing=1.4)
        y = 3.5
        for label, value in rows:
            add_text(s, x + 0.3, y, 3.5, 0.3,
                     label, size=10, bold=True, color=ACCENT, font=FONT_CN)
            add_text(s, x + 0.3, y + 0.3, 3.5, 0.4,
                     value, size=11, color=WHITE, font=FONT_CN)
            y += 0.95

    # bottom line
    add_text(s, 0.6, 6.85, 12.1, 0.3,
             "策略核心：先用「开发者免费」拿到 SDK 渗透，再把 ToB / 行业方案商业化。",
             size=11, color=GOOD, font=FONT_CN)
    add_footer(s, 14)
    add_notes(s, "GTM · 三步从免费 → 行业 → 企业的标准 SaaS 入场曲线。")


def slide_roadmap(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "07", "Roadmap · 已经做了什么、下一步做什么",
                       kicker="ROADMAP")

    cols = [
        ("已完成", "DONE",
         ["WASM 微内核 + VFS 自研",
          "包管理器与 Registry 雏形",
          "git / docker / kubectl 全量移植",
          "TypeScript / Go 编译器适配",
          "WebGPU 本地 27B 模型 (Bonsai)",
          "工作区持久化 + 启动器顺序编辑器",
          "Wagi Dog 桌面宠物"]),
        ("12–18 个月", "NEXT",
         ["公共 SDK 与开发文档开放",
          "首批行业付费客户落地",
          "插件分发体系闭环",
          "公开路演与行业试点",
          "合规与安全认证材料"]),
        ("中长期", "LATER",
         ["Agent 应用商店 (Marketplace)",
          "流程即服务 (PaaS)",
          "WASM 生态标准共建",
          "跨设备同步协议",
          "行业数字化转型升级"]),
    ]
    for i, (head, tag, items) in enumerate(cols):
        x = 0.6 + i * 4.07
        add_card(s, x, 1.7, 3.87, 4.9, accent_left=(i == 1))
        add_pill(s, x + 0.3, 1.85, tag, w=0.9, h=0.3, size=10)
        add_text(s, x + 1.3, 1.85, 2.4, 0.4,
                 head, size=18, bold=True, color=WHITE, font=FONT_CN)
        y = 2.4
        for item in items:
            add_text(s, x + 0.3, y, 0.3, 0.4,
                     "▸", size=14, bold=True, color=ACCENT)
            add_text(s, x + 0.55, y, 3.2, 0.6,
                     item, size=12, color=WHITE, font=FONT_CN, line_spacing=1.3)
            y += 0.55

    # next milestone callout
    add_card(s, 0.6, 6.75, 12.1, 0.45, fill=BG_CARD, line_color=ACCENT)
    add_text(s, 0.85, 6.78, 11.7, 0.4,
             "下一个里程碑 · NEXT MILESTONE: SDK 上线 · 首批行业付费客户 · 插件分发闭环",
             size=12, bold=True, color=ACCENT, font=FONT_CN)
    add_footer(s, 15)
    add_notes(s, "Roadmap · 已完成的事都是真事；下一步是 12-18 个月本轮融资覆盖的范围。")


def slide_team(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "08", "Team · 我们已经在这里",
                       kicker="TEAM")

    # left: founder
    add_card(s, 0.6, 1.7, 4.0, 5.0, accent_left=True)
    # avatar — use the brand mark square
    if ASSET_LOGO.exists():
        s.shapes.add_picture(str(ASSET_LOGO),
                             Inches(0.85), Inches(1.95),
                             width=Inches(1.4), height=Inches(1.4))
    add_text(s, 2.4, 1.95, 2.0, 0.4,
             "FOUNDER · 全职", size=10, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 2.4, 2.3, 2.0, 0.4,
             "邹航标", size=18, bold=True, color=WHITE, font=FONT_CN)
    add_text(s, 2.4, 2.65, 2.0, 0.3,
             "Hangbiao Zou", size=11, color=ACCENT_2, font=FONT_EN)
    add_paragraphs(s, 0.85, 3.55, 3.5, 3.0, [
        "WebAssembly / 浏览器架构深度实践者",
        "FydeOS · 国内首家 ChromiumOS 团队工程经验",
        "AI Agent 运行时与开源社区经验",
        "RustConf China / HN / 中文开发者大会演讲",
    ], size=11, color=WHITE, line_spacing=1.4, bullet=True)

    # right: post-funding plan
    add_text(s, 4.85, 1.7, 7.8, 0.4,
             "本轮融资后扩招 · 4 个核心岗位",
             size=16, bold=True, color=WHITE, font=FONT_CN)

    roles = [
        ("Web & WASM 集成开发",
         "WebAssembly / 浏览器架构方向 · 推动 SDK 公共化"),
        ("DevOps & 安全攻防",
         "FydeOS 一线工程经验 · 安全 + 云平台自动化"),
        ("Agent & 插件生态",
         "Jeff Lindsay 等国际开发者常态交流 · Registry 闭环"),
        ("产品运营 & 社区",
         "黑客松路演 · 开发者运营 · 用户支持与行业方案"),
    ]
    for i, (h, b) in enumerate(roles):
        row = i // 2
        col = i % 2
        x = 4.85 + col * 4.0
        y = 2.2 + row * 1.55
        add_card(s, x, y, 3.8, 1.4, fill=BG_PANEL)
        add_text(s, x + 0.2, y + 0.15, 3.6, 0.35,
                 h, size=14, bold=True, color=ACCENT, font=FONT_CN)
        add_text(s, x + 0.2, y + 0.55, 3.6, 0.85,
                 b, size=11, color=WHITE, font=FONT_CN, line_spacing=1.35)

    # advisors strip
    add_text(s, 4.85, 5.5, 7.8, 0.4,
             "正在接洽的外部支持",
             size=14, bold=True, color=ACCENT, font=FONT_CN)
    add_paragraphs(s, 4.85, 5.9, 7.8, 1.1, [
        "云厂商 / WebAssembly 基金会成员：标准共建与生态对接",
        "国际开源项目核心贡献者：插件与 Registry 体系合作",
        "港科大（上海）产教融合中心：孵化器 + 产业资源对接",
    ], size=11, color=WHITE, line_spacing=1.4, bullet=True)

    add_footer(s, 16)
    add_notes(s, "Team · 创始人全职 + 4 个待招岗位 + 正在接洽的外部支持。")


def slide_financials(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "09", "财务预测与核心 KPI",
                       kicker="FINANCIALS")

    # KPI table
    headers = ["指标 / Metric", "2027", "2028", "2029"]
    rows = [
        ("付费客户",     "150 家",   "600 家",    "1,800 家"),
        ("插件商品数",   "27 类",    "70 类",     "150 类+"),
        ("ARR (万元)",   "80",       "250",       "800–1,000"),
        ("续订率",       "60%+",     "65%+",      "70%+"),
        ("毛利率",       "55%",      "65%",       "70%+"),
    ]
    table_x = 0.6
    table_y = 1.7
    col_widths = [3.0, 1.4, 1.4, 1.4]
    row_h = 0.5
    # header
    cur_x = table_x
    for w, head in zip(col_widths, headers):
        add_rect(s, cur_x, table_y, w, row_h, fill=BG_PANEL, line_color=LINE)
        add_text(s, cur_x + 0.1, table_y + 0.1, w - 0.2, row_h - 0.1,
                 head, size=11, bold=True, color=ACCENT, font=FONT_EN,
                 align=PP_ALIGN.CENTER)
        cur_x += w
    # rows
    for i, row in enumerate(rows):
        y = table_y + row_h * (i + 1)
        cur_x = table_x
        for j, (w, cell) in enumerate(zip(col_widths, row)):
            fill = BG_CARD if i % 2 == 0 else BG_PANEL
            add_rect(s, cur_x, y, w, row_h, fill=fill, line_color=LINE)
            color = WHITE if j == 0 else ACCENT
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            add_text(s, cur_x + 0.1, y + 0.1, w - 0.2, row_h - 0.1,
                     cell, size=12, bold=(j != 0), color=color,
                     align=align)
            cur_x += w

    # right: big stats
    stats = [
        ("70%+", "成熟期毛利率"),
        ("2029", "预期盈亏平衡"),
        ("5–8×", "3–5 年投资回报预期"),
        ("300 万+", "第三年营收突破（万元）"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 8.1 + (i % 2) * 2.4
        y = 1.7 + (i // 2) * 1.6
        add_card(s, x, y, 2.2, 1.45, accent_left=True)
        add_text(s, x + 0.15, y + 0.2, 2.0, 0.5,
                 num, size=22, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, x + 0.15, y + 0.85, 2.0, 0.5,
                 label, size=10, color=ACCENT_2, font=FONT_CN, line_spacing=1.3)

    # bottom: assumptions
    add_text(s, 0.6, 5.05, 12, 0.4,
             "ASSUMPTIONS · 关键假设",
             size=11, bold=True, color=ACCENT, font=FONT_EN)
    add_paragraphs(s, 0.6, 5.45, 12.1, 1.5, [
        ("订阅定价", "600–1000 RMB/席位/年 · 与大陆同类 SaaS 对标处于 70-80% 分位"),
        ("客户结构演进", "2027 90% 中小团队 → 2029 50% 中大客户（行业头部带动复购）"),
        ("分成收入", "插件 / Agent 分成 15-30% · 与 HubSpot / Slack App Marketplace 模型对齐"),
        ("成本结构", "CDN + 推理 GPU + 人力 · 12 个月内固定成本可控，毛利率随规模抬升"),
    ], size=11, color=WHITE, line_spacing=1.4, bullet=True)

    add_footer(s, 17)
    add_notes(s, "Financials · KPI 表格 + 四大数字 + 假设。避免承诺过满的曲线图，只给关键节点。")


def slide_funding(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s)
    add_section_header(s, "10", "融资计划：300 万 RMB · 18 个月",
                       kicker="FUNDING")

    # big ask number
    add_card(s, 0.6, 1.7, 12.1, 1.5, accent_left=True)
    add_text(s, 0.95, 1.85, 5, 0.4,
             "ROUND 1 ASKING", size=10, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 0.95, 2.2, 5, 0.7,
             "300 万 RMB", size=36, bold=True, color=WHITE, font=FONT_EN)
    add_text(s, 0.95, 2.85, 5, 0.3,
             "覆盖 18 个月研发 / 市场 / 团队扩展周期", size=12, color=ACCENT_2, font=FONT_CN)

    add_text(s, 7.0, 1.85, 5.5, 0.4,
             "RUNWAY", size=10, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 7.0, 2.2, 5.5, 0.7,
             "18 个月 · 月均 15–20 万", size=22, bold=True, color=WHITE, font=FONT_EN)
    add_text(s, 7.0, 2.85, 5.5, 0.3,
             "目标 12 个月内达成 PMF 信号 · 18 个月内进入 Series A 准备期",
             size=12, color=ACCENT_2, font=FONT_CN)

    # use of funds
    add_text(s, 0.6, 3.4, 12, 0.4,
             "USE OF FUNDS · 资金用途",
             size=11, bold=True, color=ACCENT, font=FONT_EN)
    funds = [
        ("60%", "核心研发与产品完善",  "WASM 微内核 · VFS · Registry 平台 · 多语言 SDK 开放"),
        ("20%", "市场开拓与生态合作", "行业试点 DEMO · 峰会参与 · ISV 与云市场集成"),
        ("15%", "团队扩充与人才引进", "WASM / AI Agent / 前端安全研发骨干 + 产品运营"),
        ("5%",  "运营与基础设施",       "CDN · 演示环境 · 基础办公与合规开销"),
    ]
    bar_x = 0.6
    bar_y = 3.85
    bar_total = 12.1
    cum = 0
    bar_colors = [ACCENT, RGBColor(0x6E, 0xB4, 0xFF), RGBColor(0xA6, 0xC8, 0xFF), MUTED]
    for i, (pct, head, body) in enumerate(funds):
        seg = bar_total * (int(pct[:-1]) / 100)
        add_rect(s, bar_x + cum, bar_y, seg, 0.4,
                 fill=bar_colors[i], line_color=BG_DEEP)
        # label under segment
        add_text(s, bar_x + cum + 0.1, bar_y + 0.05, seg - 0.2, 0.3,
                 pct, size=12, bold=True, color=BG_DEEP, font=FONT_EN)
        cum += seg

    # detail rows
    y = 4.45
    for pct, head, body in funds:
        add_card(s, 0.6, y, 12.1, 0.55, fill=BG_PANEL)
        add_text(s, 0.85, y + 0.1, 1.0, 0.4,
                 pct, size=13, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, 2.0, y + 0.1, 3.5, 0.4,
                 head, size=11, bold=True, color=WHITE, font=FONT_CN)
        add_text(s, 5.6, y + 0.13, 7.0, 0.4,
                 body, size=10, color=ACCENT_2, font=FONT_CN)
        y += 0.62

    add_footer(s, 18)
    add_notes(s, "Funding · 一眼看清金额 / runway / 资金分配。避免堆数字。")


def slide_ask(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    set_slide_bg(s, gradient_to=BG_PANEL)
    add_text(s, 0.6, 0.7, 12, 0.5, "THE ASK",
             size=12, bold=True, color=ACCENT, font=FONT_EN)
    add_text(s, 0.6, 1.05, 12, 0.7,
             "Why Now · Why Us · What We Need",
             size=28, bold=True, color=WHITE, font=FONT_EN)

    # three columns: why now / why us / ask
    cols = [
        ("WHY NOW",
         "浏览器正在成为新 OS",
         ["WebAssembly 开放标准成熟，多语言 sandbox 落地",
          "WebGPU 让 27B 本地模型首次在浏览器内跑通",
          "Agent 厂商分发成本高企，「装不动」成为漏斗最大瓶颈",
          "GDPR / 欧盟 AI 法案 / 中国算法备案推动隔离层必备"]),
        ("WHY US",
         "唯一全栈经验",
         ["FydeOS · 国内首家 ChromiumOS 团队工程底子",
          "GearShell 自举验证：浏览器内构建浏览器内",
          "Bonsai 27B · Wagi Dog · Browser.js 已构成产品矩阵",
          "国际开发者社区认可（HN / RustConf / Jeff Lindsay）"]),
        ("WHAT WE NEED",
         "本轮要什么",
         ["300 万 RMB · 18 个月 runway",
          "公共 SDK 上线 · 首批付费客户落地 · 插件分发闭环",
          "港科大（上海）产教融合中心孵化器入驻与产业对接",
          "投资人：技术 + 开发者生态 + ToB 渠道资源优先"]),
    ]
    for i, (tag, head, items) in enumerate(cols):
        x = 0.6 + i * 4.07
        add_card(s, x, 1.9, 3.87, 4.4, accent_left=True)
        add_text(s, x + 0.3, 2.05, 3.5, 0.3,
                 tag, size=11, bold=True, color=ACCENT, font=FONT_EN)
        add_text(s, x + 0.3, 2.4, 3.5, 0.5,
                 head, size=18, bold=True, color=WHITE, font=FONT_CN)
        y = 3.0
        for it in items:
            add_text(s, x + 0.3, y, 0.3, 0.4,
                     "▸", size=14, bold=True, color=ACCENT)
            add_text(s, x + 0.55, y, 3.2, 0.7,
                     it, size=11, color=WHITE, font=FONT_CN, line_spacing=1.35)
            y += 0.65

    # closing line
    add_card(s, 0.6, 6.4, 12.1, 0.7, accent_left=True, line_color=ACCENT)
    add_text(s, 0.95, 6.55, 11.7, 0.4,
             "期待与您共建 Agent 分发生态 —— 让 AI Agent 像网页一样开箱即用。",
             size=14, bold=True, color=WHITE, font=FONT_CN)
    add_text(s, 11.0, 7.1, 2, 0.3, "Round 1 · 19 / 20",
             size=9, color=MUTED, font=FONT_EN, align=PP_ALIGN.RIGHT)
    add_notes(s, "Ask · 收尾页。三栏分别回答 why now / why us / ask，最后一句 CTA。")


# ---------- assembly ----------

def main():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)

    builders = [
        slide_cover,
        slide_agenda,
        slide_tldr,
        slide_problem,
        slide_voice_from_field,
        slide_solution,
        slide_demo,
        slide_architecture,
        slide_why_us,
        slide_traction,
        slide_market,
        slide_customers,
        slide_business,
        slide_gtm,
        slide_roadmap,
        slide_team,
        slide_financials,
        slide_funding,
        slide_ask,
    ]
    for b in builders:
        b(p)

    out = WORKSPACE / "神奇海螺队 - Round 1 Pitch Deck.pptx"
    p.save(out)
    size_mb = out.stat().st_size / 1024 / 1024
    print(f"Saved {out}")
    print(f"Slides: {len(p.slides)}")
    print(f"Size:   {size_mb:.2f} MB")


if __name__ == "__main__":
    main()