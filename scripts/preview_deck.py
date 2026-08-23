#!/usr/bin/env python3
"""Render the built PPT as a static HTML preview (no external deps).

This mirrors the slides' key text blocks + layout at 13.333 x 7.5 in
so we can spot-check positioning visually without Keynote.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

WORKSPACE = Path("/Users/gear/Documents/GitHub/gearshell")
PPTX = WORKSPACE / "神奇海螺队 - Round 1 Pitch Deck.pptx"
OUT = WORKSPACE / "scripts" / "deck_preview.html"

EMU_PER_INCH = 914400


def emu_to_in(v):
    if v is None:
        return 0
    return v / EMU_PER_INCH


def shape_kind(sh):
    if sh.shape_type == 13:  # PICTURE
        return "picture"
    if sh.has_text_frame:
        return "text"
    return "shape"


def text_summary(tf):
    lines = []
    for p in tf.paragraphs:
        for r in p.runs:
            t = r.text.strip()
            if t:
                lines.append(t)
    return lines


def main():
    p = Presentation(PPTX)
    sw, sh = p.slide_width / EMU_PER_INCH, p.slide_height / EMU_PER_INCH
    html = [f"<!doctype html><meta charset='utf-8'>",
            f"<title>Deck Preview · {PPTX.name}</title>",
            "<style>",
            "body{margin:0;background:#04060d;color:#fff;font:14px -apple-system,'PingFang SC','Microsoft YaHei',sans-serif}",
            ".grid{display:grid;grid-template-columns:repeat(2,1fr);gap:24px;padding:24px}",
            ".slide{position:relative;background:#0b1220;border:1px solid #2b3a67;border-radius:12px;aspect-ratio:16/9;overflow:hidden}",
            ".slide .label{position:absolute;top:8px;left:8px;background:#8ab4ff;color:#0b1220;padding:2px 8px;border-radius:6px;font-size:11px;font-weight:700}",
            ".shape{position:absolute;box-sizing:border-box}",
            ".shape.text{padding:4px 8px;line-height:1.2;overflow:hidden}",
            ".shape.shape{background:#1c2b54;border:1px solid #2b3a67;border-radius:6px}",
            ".shape.accent{border-left:3px solid #8ab4ff}",
            ".shape.picture{background:repeating-linear-gradient(45deg,#1c2b54,#1c2b54 6px,#14213d 6px,#14213d 12px);border:1px dashed #8ab4ff;color:#8ab4ff;font-size:11px;display:flex;align-items:center;justify-content:center}",
            "</style>",
            f"<h1 style='padding:24px'>{PPTX.name} · {len(p.slides)} slides · {sw:.2f} x {sh:.2f} in</h1>",
            "<div class='grid'>"]

    for i, s in enumerate(p.slides, 1):
        html.append(f"<div class='slide'><span class='label'>{i:02d}</span>")
        # shapes in z-order
        for shp in s.shapes:
            l = emu_to_in(shp.left) / sw * 100
            t = emu_to_in(shp.top) / sh * 100
            w = emu_to_in(shp.width) / sw * 100
            h = emu_to_in(shp.height) / sh * 100
            style = f"left:{l:.2f}%;top:{t:.2f}%;width:{w:.2f}%;height:{h:.2f}%;"
            kind = shape_kind(shp)
            if kind == "text":
                lines = text_summary(shp.text_frame)
                if not lines:
                    continue
                # detect accent-left card
                accent_class = ""
                if shp.fill.type and shp.fill.type == 1:
                    rgb = shp.fill.fore_color.rgb
                    if rgb is not None and str(rgb) == "1C2B54":
                        accent_class = " accent"
                txt = " · ".join(lines)
                # shrink long text
                fs = "10px" if len(txt) > 80 else "12px"
                if len(txt) > 200:
                    fs = "8px"
                html.append(f"<div class='shape text{accent_class}' style='{style}font-size:{fs}'>{txt}</div>")
            elif kind == "picture":
                html.append(f"<div class='shape picture' style='{style}'>IMG</div>")
            else:
                html.append(f"<div class='shape shape' style='{style}'></div>")
        html.append("</div>")
    html.append("</div></body>")
    OUT.write_text("\n".join(html))
    print(f"Wrote {OUT}")
    print(f"Open: file://{OUT}")


if __name__ == "__main__":
    main()